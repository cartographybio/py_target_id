__all__ = ['Slides']

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class Slides:
    def __init__(self):
        self.prs = None
        self.filename = None
        self.footer_left = "Confidential | Do Not Distribute"
        self.footer_right = "2025 | Cartography Biosciences"
        self.footer_logo = "cartography_footer.png"
        self.title_color = RGBColor(0, 51, 102)  # Dark blue
        self.footer_color = RGBColor(166, 166, 166)  # Gray
        self.text_margin = Inches(0.3)
        self.title_height = Inches(0.7)
        self.title_margin = Inches(0.15)
    
    def start(self, filename, footer_left=None, footer_right=None, footer_logo=None):
        """Initialize a new PowerPoint presentation"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # Widescreen 16:9
        self.prs.slide_height = Inches(7.5)
        self.filename = filename
        if footer_left:
            self.footer_left = footer_left
        if footer_right:
            self.footer_right = footer_right
        if footer_logo:
            self.footer_logo = footer_logo
        return self
    
    def _add_footer(self, slide, slide_width, slide_height):
        """Add footer to slide with left and right text and logo"""
        footer_height = Inches(0.3)
        footer_top = slide_height - footer_height
        footer_font_size = 8
        
        # Add clean gray line above footer (no shadow or border)
        footer_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.text_margin,
            footer_top - Pt(1),
            slide_width - 2 * self.text_margin,
            Pt(0.5)  # Slightly thicker for visibility
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self.footer_color
        footer_line.shadow.inherit = False  # Remove shadow
        footer_line.line.fill.background()  # Remove outline color
        footer_line.line.width = Pt(0)      # No outline
        
        # Get slide number
        slide_number = len(self.prs.slides)
        
        # Add logo in bottom left
        try:
            logo = slide.shapes.add_picture(
                self.footer_logo,
                self.text_margin,
                footer_top + Inches(0.05),
                height=footer_height * 0.5 # 80% of footer height
            )
            logo_right = logo.left + logo.width
        except:
            logo_right = self.text_margin  # Skip logo if file not found
        
        # Left footer text (to the right of logo with minimal spacing)
        text_spacing = Inches(0.1)
        left_footer = slide.shapes.add_textbox(
            logo_right + text_spacing,
            footer_top,
            slide_width / 2 - logo_right - text_spacing,
            footer_height
        )
        left_frame = left_footer.text_frame
        left_frame.vertical_anchor = 1
        left_p = left_frame.paragraphs[0]
        left_p.text = self.footer_left
        left_p.font.size = Pt(footer_font_size)
        left_p.font.name = 'Arial'
        left_p.font.color.rgb = self.footer_color
        left_p.alignment = PP_ALIGN.LEFT
        
        # Right footer with page number
        right_footer = slide.shapes.add_textbox(
            slide_width / 2,
            footer_top,
            slide_width / 2 - self.text_margin,
            footer_height
        )
        right_frame = right_footer.text_frame
        right_frame.vertical_anchor = 1
        right_p = right_frame.paragraphs[0]
        right_p.text = f"{self.footer_right} | {slide_number}"
        right_p.font.size = Pt(footer_font_size)
        right_p.font.name = 'Arial'
        right_p.font.color.rgb = self.footer_color
        right_p.alignment = PP_ALIGN.RIGHT
    
    def add_slide_text_plot(self, ratio, plot, text, 
                           title="Title",
                           text_font_size=18, 
                           title_font_size=32,
                           text_alignment="left",
                           plot_on_right=True):
        """Add a slide with text and plot in horizontal layout"""
        text_ratio, plot_ratio = map(int, ratio.split(":"))
        total = text_ratio + plot_ratio
        text_width_ratio = text_ratio / total
        plot_width_ratio = plot_ratio / total
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        content_top = self.text_margin + self.title_height + self.title_margin
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        usable_width = slide_width - self.text_margin
        usable_height = slide_height - content_top
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.text_margin,
            self.text_margin,
            slide_width - self.text_margin,
            self.title_height
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.name = 'Arial'
        title_p.font.bold = True
        title_p.font.color.rgb = self.title_color
        title_p.alignment = PP_ALIGN.LEFT
        
        # --- FIXED BLUE LINE (no gray shadow) ---
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.text_margin,
            self.text_margin + self.title_height,
            slide_width - 2 * self.text_margin,
            Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.title_color
        line.shadow.inherit = False
        line.line.fill.background()
        line.line.width = Pt(0)
        
        # Footer
        self._add_footer(slide, slide_width, slide_height)
        
        # Layout
        text_width = usable_width * text_width_ratio
        plot_width = usable_width * plot_width_ratio
        
        if plot_on_right:
            text_left = self.text_margin
            plot_left = self.text_margin + text_width
            plot_width_actual = slide_width - plot_left
        else:
            plot_left = 0
            plot_width_actual = usable_width * plot_width_ratio
            text_left = plot_left + plot_width_actual
        
        textbox = slide.shapes.add_textbox(text_left, content_top, text_width, usable_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_font_size)
        p.font.name = 'Arial'
        alignment_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = alignment_map.get(text_alignment.lower(), PP_ALIGN.LEFT)
        
        pic = slide.shapes.add_picture(plot, plot_left, content_top)
        orig_width, orig_height = pic.width, pic.height
        aspect_ratio = orig_width / orig_height
        max_width, max_height = plot_width_actual, usable_height
        
        if orig_width / max_width > orig_height / max_height:
            new_width = max_width
            new_height = max_width / aspect_ratio
        else:
            new_height = max_height
            new_width = max_height * aspect_ratio
        
        pic.width = int(new_width)
        pic.height = int(new_height)
        pic.left = slide_width - pic.width if plot_on_right else plot_left
        pic.top = content_top + int((usable_height - new_height) / 2)
        
        self._add_footer(slide, slide_width, slide_height)
        return self
    
    def add_slide_full_plot(self, plot, title="", title_font_size=32):
        """Add a slide with just a full-size plot, preserving aspect ratio"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        text_margin = Inches(0.3)
        title_height = Inches(0.7) if title else Inches(0)
        title_margin = Inches(0.15) if title else Inches(0)
        
        if title:
            title_box = slide.shapes.add_textbox(
                text_margin,
                text_margin,
                self.prs.slide_width - text_margin,
                title_height
            )
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = title
            title_p.font.size = Pt(title_font_size)
            title_p.font.name = 'Arial'
            title_p.font.bold = True
            title_p.alignment = PP_ALIGN.LEFT
        
        content_top = text_margin + title_height + title_margin
        pic = slide.shapes.add_picture(plot, 0, content_top)
        orig_width, orig_height = pic.width, pic.height
        aspect_ratio = orig_width / orig_height
        max_width = self.prs.slide_width
        max_height = self.prs.slide_height - content_top
        
        if orig_width / max_width > orig_height / max_height:
            pic.width = int(max_width)
            pic.height = int(max_width / aspect_ratio)
        else:
            pic.height = int(max_height)
            pic.width = int(max_height * aspect_ratio)
        
        pic.left = int((max_width - pic.width) / 2)
        pic.top = content_top + int((max_height - pic.height) / 2)
        return self
    
    def add_slide_plot_subtitle(self, plot, title="Title", subtitle="", 
                                title_height_ratio=0.15, 
                                plot_height_ratio=0.70, 
                                subtitle_height_ratio=0.15,
                                title_font_size=32,
                                subtitle_font_size=14):
        """Add a slide with title, plot, and subtitle in vertical layout"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        available_height = slide_height - self.text_margin
        title_height = available_height * title_height_ratio
        plot_height = available_height * plot_height_ratio
        subtitle_height = available_height * subtitle_height_ratio
        title_top = self.text_margin
        plot_top = self.text_margin + self.title_height + self.title_margin
        subtitle_top = self.text_margin + self.title_height + self.title_margin + plot_height
        
        # Title
        title_box = slide.shapes.add_textbox(
            self.text_margin,
            title_top,
            slide_width - self.text_margin,
            self.title_height
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.name = 'Arial'
        title_p.font.bold = True
        title_p.font.color.rgb = self.title_color
        title_p.alignment = PP_ALIGN.LEFT
        
        # --- FIXED BLUE LINE (no gray shadow) ---
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.text_margin,
            self.text_margin + self.title_height,
            slide_width - 2 * self.text_margin,
            Pt(1.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.title_color
        line.shadow.inherit = False
        line.line.fill.background()
        line.line.width = Pt(0)
        
        self._add_footer(slide, slide_width, slide_height)
        
        pic = slide.shapes.add_picture(plot, 0, plot_top)
        orig_width, orig_height = pic.width, pic.height
        aspect_ratio = orig_width / orig_height
        max_width = slide_width
        max_height = plot_height
        
        if orig_width / max_width > orig_height / max_height:
            pic.width = int(max_width)
            pic.height = int(max_width / aspect_ratio)
        else:
            pic.height = int(max_height)
            pic.width = int(max_height * aspect_ratio)
        
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = int(plot_top + (plot_height - pic.height) / 2)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.text_margin,
                subtitle_top,
                slide_width - self.text_margin,
                subtitle_height
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_frame.vertical_anchor = 1
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(subtitle_font_size)
            subtitle_p.font.name = 'Arial'
            subtitle_p.alignment = PP_ALIGN.LEFT
        
        return self
    
    def add_slide_title(self, title, subtitle=""):
        """Add a title slide with gradient background"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # Add gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45.0
        fill.gradient_stops[0].color.rgb = RGBColor(20, 90, 130)  # Teal blue
        fill.gradient_stops[1].color.rgb = RGBColor(80, 50, 120)  # Purple
        
        # Add title - large, white, bold, left-aligned, upper-left area
        title_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        title_p = title_frame.paragraphs[0]
        title_p.text = title.upper()
        title_p.font.size = Pt(72)
        title_p.font.name = 'Arial'
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)  # White
        title_p.alignment = PP_ALIGN.LEFT
        
        # Add subtitle if provided - smaller, light blue, below title
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(4.5),
                Inches(8),
                Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle.upper()
            subtitle_p.font.size = Pt(36)
            subtitle_p.font.name = 'Arial'
            subtitle_p.font.color.rgb = RGBColor(100, 200, 230)  # Light blue
            subtitle_p.alignment = PP_ALIGN.LEFT
        
        # Add Cartography Biosciences logo/text in bottom right
        logo_box = slide.shapes.add_textbox(
            slide_width - Inches(4),
            slide_height - Inches(1.5),
            Inches(3.5),
            Inches(1)
        )
        logo_frame = logo_box.text_frame
        logo_frame.word_wrap = False
        
        # First line: CARTOGRAPHY
        logo_p1 = logo_frame.paragraphs[0]
        logo_p1.text = "CARTOGRAPHY"
        logo_p1.font.size = Pt(28)
        logo_p1.font.name = 'Arial'
        logo_p1.font.bold = True
        logo_p1.font.color.rgb = RGBColor(255, 255, 255)
        logo_p1.alignment = PP_ALIGN.RIGHT
        
        # Second line: BIOSCIENCES
        logo_p2 = logo_frame.add_paragraph()
        logo_p2.text = "BIOSCIENCES"
        logo_p2.font.size = Pt(18)
        logo_p2.font.name = 'Arial'
        logo_p2.font.color.rgb = RGBColor(255, 255, 255)
        logo_p2.alignment = PP_ALIGN.RIGHT
        
        # Add footer text in bottom left
        footer_box = slide.shapes.add_textbox(
            self.text_margin,
            slide_height - Inches(0.5),
            Inches(4),
            Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        footer_p = footer_frame.paragraphs[0]
        footer_p.text = self.footer_left
        footer_p.font.size = Pt(10)
        footer_p.font.name = 'Arial'
        footer_p.font.color.rgb = RGBColor(150, 180, 200)  # Light blue-gray
        footer_p.alignment = PP_ALIGN.LEFT
        
        return self

    def add_slide_section(self, title, subtitle=""):
        """Add a section divider slide with gradient background"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # Add gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135.0  # Different angle for variety
        fill.gradient_stops[0].color.rgb = RGBColor(50, 100, 150)  # Blue
        fill.gradient_stops[1].color.rgb = RGBColor(70, 130, 180)  # Lighter blue
        
        # Add title - large, white, bold, right-aligned, upper-right area
        title_box = slide.shapes.add_textbox(
            slide_width - Inches(7),
            Inches(1.5),
            Inches(6),
            Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(60)
        title_p.font.name = 'Arial'
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)  # White
        title_p.alignment = PP_ALIGN.RIGHT
        
        # Add subtitle if provided - smaller, light, right-aligned, center area
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                slide_width - Inches(7),
                slide_height / 2 + Inches(0.5),
                Inches(6),
                Inches(1.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(32)
            subtitle_p.font.name = 'Arial'
            subtitle_p.font.color.rgb = RGBColor(200, 220, 240)  # Light blue
            subtitle_p.alignment = PP_ALIGN.RIGHT
        
        # Add footer
        self._add_footer(slide, slide_width, slide_height)
        
    def add_slide_blank(self):
        """Add a blank slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        return self
    
    def close(self):
        """Save and close the PowerPoint presentation"""
        if self.prs and self.filename:
            self.prs.save(self.filename)
            print(f"Presentation saved as '{self.filename}'")
        else:
            print("No presentation to save")
